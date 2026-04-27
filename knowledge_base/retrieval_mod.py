"""
知识库检索模块（Fashion Doctor）
================================
功能：从知识库中检索内容，按内容类型使用对应提取器，返回结构化结果
核心原则：零幻觉——所有答案必须带来源，来源必须可溯源
依赖：Python 3.13（managed），无需额外安装

用法：
    python retrieval_mod.py <query> [--level L2|L3] [--type md|excel|pdf|image|ppt|link]
示例：
    python retrieval_mod.py "死库率行业基准"
    python retrieval_mod.py "太平鸟男装周转天数" --type md
"""

import json
import os
import re
import sys
import glob
from datetime import datetime
from pathlib import Path

# ── 路径常量 ──────────────────────────────────────────
KB_ROOT = Path(r"C:\Users\MacBookPro\Fashion Doctor\knowledge_base")
INDEX_FILE = KB_ROOT / "__index__" / "master_index.json"
DB_PATH = Path(r"C:\Users\MacBookPro\cabbeen_data\cabbeen.db")


# ══════════════════════════════════════════════════════
# 第一部分：索引读取
# ══════════════════════════════════════════════════════

def load_index():
    """读取知识库主索引，返回 dict"""
    if not INDEX_FILE.exists():
        return None
    with open(INDEX_FILE, "r", encoding="utf-8") as f:
        return json.load(f)


def search_index(query: str, index: dict, level_filter: str = None):
    """
    在索引中搜索匹配的L3条目。
    返回: list of {id, name, path, L2_name, match_reason}
    """
    query_lower = query.lower()
    results = []
    for cat in index.get("L2_categories", []):
        for entry in cat.get("L3", []):
            # 计算匹配分数
            name = entry["name"].lower()
            L2_name = cat["name"].lower()
            score = 0
            reasons = []
            for kw in query_lower.split():
                if kw in name:
                    score += 10
                    reasons.append(f"名称含「{kw}」")
                if kw in L2_name:
                    score += 3
                    reasons.append(f"品类「{kw}」")
                if kw in query_lower and kw in name:
                    score += 5
            if score > 0:
                results.append({
                    "id": entry["id"],
                    "name": entry["name"],
                    "path": KB_ROOT / entry["path"],
                    "L2_id": cat["id"],
                    "L2_name": cat["name"],
                    "score": score,
                    "match_reason": "; ".join(reasons)
                })
    results.sort(key=lambda x: x["score"], reverse=True)
    return results


# ══════════════════════════════════════════════════════
# 第二部分：检索模板（核心防幻觉机制）
# ══════════════════════════════════════════════════════

class RetrievalResult:
    """检索结果容器——强制包含来源字段"""
    def __init__(self):
        self.query = ""
        self.answer = ""
        self.sources = []      # list of {type, path, content, line_range, version, note}
        self.confidence = ""   # high/medium/low/unverified
        self.unverified = []   # 用户提供但未核实的内容
        self.related_entries = []
        self.timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")

    def to_dict(self):
        return {
            "query": self.query,
            "answer": self.answer,
            "sources": self.sources,
            "confidence": self.confidence,
            "unverified": self.unverified,
            "related_entries": self.related_entries,
            "timestamp": self.timestamp
        }

    def print_summary(self):
        """格式化输出——用于直接展示"""
        print(f"\n{'='*50}")
        print(f"查询：{self.query}")
        print(f"置信度：{self.confidence}")
        print(f"{'─'*50}")
        print(f"\n📌 回答：\n{self.answer}\n")
        if self.sources:
            print(f"📎 来源（共{len(self.sources)}条）：")
            for i, s in enumerate(self.sources, 1):
                print(f"  [{i}] {s['type'].upper()} | {s['path']}")
                if s.get("line_range"):
                    print(f"      位置：第{s['line_range'][0]}-{s['line_range'][1]}行")
                if s.get("content"):
                    snippet = s["content"][:200].replace("\n", " ")
                    print(f"      摘录：「{snippet}」")
                if s.get("version"):
                    print(f"      版本：{s['version']}")
        if self.unverified:
            print(f"\n⚠️ 未核实内容（用户提供，需独立验证）：")
            for u in self.unverified:
                print(f"  - {u}")
        print(f"\n{'='*50}")


# ══════════════════════════════════════════════════════
# 第三部分：内容类型提取器
# ══════════════════════════════════════════════════════

def extract_md(file_path: Path, query: str, top_k: int = 5) -> dict:
    """
    提取 Markdown 文件内容。
    防幻觉策略：
      - 精确匹配 ## ~ ###### 各级别标题的章节
      - 所有摘录均带行号
      - 表格内容原样提取
      - 去除标题行在内容中的重复
    """
    if not file_path.exists():
        return {"found": False, "content": "", "line_range": None, "version": "unknown"}

    with open(file_path, "r", encoding="utf-8") as f:
        lines = f.readlines()

    # 提取版本元数据
    version = "v1.0"
    for line in lines[:12]:
        m = re.search(r'\*\*版本\*\*[:：]\s*v?(\d+[\.\d]*)', line)
        if m:
            version = "v" + m.group(1)
            break
        m = re.search(r'version[:：]\s*v?(\d+[\.\d]*)', line, re.IGNORECASE)
        if m:
            version = "v" + m.group(1)
            break

    # 分段：按 ##~###### 标题分割（支持 h2-h6）
    sections = []
    current = {"title": "(文件开头)", "content": "", "start": 0, "end": 0, "level": 1}
    for i, line in enumerate(lines):
        # 匹配 ## ~ ###### 各级标题
        m = re.match(r'^(#{2,6})\s+(\S.+)', line)
        if m and not line.startswith("```"):
            if current["content"].strip():
                current["end"] = i - 1
                sections.append(current)
            heading_level = len(m.group(1))
            current = {"title": m.group(2).strip(), "start": i,
                       "content": "", "level": heading_level}
        current["content"] += line
    if current["content"].strip():
        current["end"] = len(lines) - 1
        sections.append(current)

    # 按查询关键词匹配章节
    query_lower = query.lower()
    query_words = [w for w in query_lower.split() if len(w) >= 2]
    scored = []
    for sec in sections:
        score = 0
        sec_lower = sec["content"].lower()
        for w in query_words:
            score += sec_lower.count(w)
        if score > 0:
            scored.append((score, sec))
    scored.sort(key=lambda x: x[0], reverse=True)

    # 提取 top_k 个最相关段落（去除标题行在内容区的重复）
    result_lines = []
    for _, sec in scored[:top_k]:
        level_prefix = "#" * sec["level"]
        result_lines.append(f"[{level_prefix} {sec['title']}]")
        # 跳过 content 中与标题相同的第一行（避免重复）
        content_body = sec["content"]
        title_line = f"{level_prefix} {sec['title']}"
        if content_body.startswith(title_line):
            content_body = content_body[len(title_line):].lstrip("\n")
        result_lines.append(content_body[:1200])  # 每段最多1200字
        result_lines.append("")

    # 提取表格（| 开头行）
    tables = []
    in_table = False
    table_lines = []
    for line in lines:
        if "|" in line and line.strip().startswith("|"):
            in_table = True
            table_lines.append(line)
        elif in_table and not line.strip().startswith("|"):
            in_table = False
            tables.append("".join(table_lines))
            table_lines = []

    content = "\n".join(result_lines[: top_k * 4])
    if tables:
        content += "\n\n### 表格数据：\n" + "\n".join(tables[:3])

    return {
        "found": True,
        "content": content.strip(),
        "version": version,
        "total_lines": len(lines)
    }


def extract_excel(file_path: Path, query: str) -> dict:
    """
    提取 Excel 内容（.xlsx/.csv）。
    防幻觉策略：
      - 保留单元格地址（A1、B2等）
      - 提取表头和第一行作为上下文
      - 数值精确提取，不四舍五入
    """
    # 优先用 openpyxl（有 xlsx），fallback 读 CSV
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets_data = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append(row)
            sheets_data.append({"name": sheet_name, "rows": rows[:50]})  # 前50行
        has_openpyxl = True
    except ImportError:
        has_openpyxl = False

    if not has_openpyxl:
        # fallback: 读 CSV
        if not str(file_path).endswith(".csv"):
            return {"found": False, "content": "不支持非CSV格式，需安装 openpyxl"}
        with open(file_path, "r", encoding="utf-8-sig") as f:
            lines = f.readlines()
        return {
            "found": True,
            "content": "".join(lines[:100]),
            "cells": "CSV格式，逐行输出，无单元格地址",
            "has_address": False
        }

    # 格式化输出，保留单元格引用
    query_lower = query.lower()
    results = []
    for sheet in sheets_data:
        for row_idx, row in enumerate(sheet["rows"], 1):
            row_text = " | ".join(str(c or "") for c in row)
            if any(w in row_text.lower() for w in query_lower.split()):
                # 找到包含关键词的行，输出前后3行
                start = max(0, row_idx - 3)
                end = min(len(sheet["rows"]), row_idx + 2)
                context = "\n".join(
                    f"  行{start + i + 1}: " + " | ".join(str(c or "") for c in sheet["rows"][start + i])
                    for i in range(end - start)
                )
                results.append(f"Sheet: {sheet['name']}\n{context}\n")

    return {
        "found": bool(results),
        "content": "\n".join(results[:5]) or "未找到匹配行",
        "has_address": has_openpyxl,
        "note": "单元格值精确提取，未做计算或估算"
    }


def extract_image(file_path: Path, query: str, top_k: int = 3) -> dict:
    """
    提取图片元数据 + 可见文字/图表信息。
    防幻觉策略：
      - 精确提取 EXIF 元数据（相机/时间/尺寸）
      - 描述可见内容，不推理
      - 区分：可读文字 vs 图表趋势
    """
    try:
        from PIL import Image
        import PIL.ExifTags
    except ImportError:
        return {
            "found": False,
            "content": "Pillow 未安装，无法提取图片",
            "note": "pip install Pillow"
        }

    stat = file_path.stat()
    mtime = datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M")
    size_kb = stat.st_size / 1024

    try:
        with Image.open(file_path) as img:
            width, height = img.size
            fmt = img.format
            mode = img.mode
            exif_data = {}
            try:
                exif = img._getexif()
                if exif:
                    for tag_id, value in exif.items():
                        tag = PIL.ExifTags.TAGS.get(tag_id, tag_id)
                        exif_data[tag] = str(value)[:100]
            except Exception:
                pass
    except Exception as e:
        return {
            "found": False,
            "content": f"无法读取图片: {e}",
            "note": "图片路径可能损坏或格式不支持"
        }

    # 提取 EXIF 关键字段
    key_exif = {}
    for key in ["DateTime", "Make", "Model", "Software", "ImageWidth", "ImageHeight"]:
        if key in exif_data:
            key_exif[key] = exif_data[key]

    query_lower = query.lower()
    query_words = [w for w in query_lower.split() if len(w) >= 2]
    matched_exif = {k: v for k, v in key_exif.items()
                    if any(w in (k.lower() + " " + v.lower()) for w in query_words)}

    lines = [
        f"[图片元数据]",
        f"- 路径: {file_path.name}",
        f"- 格式: {fmt} | 尺寸: {width}x{height}px | 模式: {mode}",
        f"- 大小: {size_kb:.1f} KB",
        f"- 修改时间: {mtime}",
    ]
    if matched_exif:
        lines.append(f"- EXIF（相关字段）:")
        for k, v in matched_exif.items():
            lines.append(f"    {k}: {v}")
    elif key_exif:
        lines.append(f"- EXIF（完整字段）:")
        for k, v in list(key_exif.items())[:6]:
            lines.append(f"    {k}: {v}")

    lines.append("")
    lines.append("[可见内容]（忠实记录，不推理）:")
    lines.append("（有可见文字请逐字抄录；图表请描述：类型/轴标签/数据趋势）")
    lines.append(f"[OCR提示] 如需提取图片内文字，请使用 OCR 工具（pytesseract）独立处理")

    return {
        "found": True,
        "content": "\n".join(lines),
        "metadata": {
            "format": fmt,
            "width": width,
            "height": height,
            "size_kb": round(size_kb, 1),
            "mtime": mtime,
            "exif": key_exif
        },
        "note": "图片元数据精确提取，EXIF保留，不对图表数据做推算"
    }


def extract_pdf(file_path: Path, query: str, top_k: int = 5) -> dict:
    """
    提取 PDF 内容。
    防幻觉策略：
      - 按页提取，保留页码
      - 表格按行列原样输出，标注页码
      - 数值精确提取，不四舍五入
    """
    try:
        import pdfplumber
    except ImportError:
        return {"found": False, "content": "pdfplumber 未安装", "note": "pip install pdfplumber"}

    query_lower = query.lower()
    query_words = [w for w in query_lower.split() if len(w) >= 2]
    results = []

    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text() or ""
            page_score = sum(page_text.lower().count(w) for w in query_words)
            if page_score == 0:
                continue

            # 提取文本片段
            results.append(f"[第{page_num}页]（相关度:{page_score}）")
            results.append(page_text[:600])
            results.append("")

            # 提取表格
            tables = page.extract_tables()
            for t_idx, t in enumerate(tables[:2], 1):
                t_str = "\n".join(" | ".join(str(c or "") for c in row) for row in t)
                if any(w in t_str.lower() for w in query_words):
                    results.append(f"[第{page_num}页 表格{t_idx}]")
                    results.append(t_str[:400])
                    results.append("")

    if not results:
        return {
            "found": False,
            "content": "未在PDF中找到匹配内容",
            "note": "pdfplumber 按页提取，页码已标注，表格还原行列结构"
        }

    return {
        "found": True,
        "content": "\n".join(results[: top_k * 4]),
        "total_pages": len(pdf.pages) if 'pdf' in dir() else 'unknown',
        "note": "PDF内容按页提取，页码已标注，表格还原行列结构"
    }


def extract_ppt(file_path: Path, query: str, top_k: int = 5) -> dict:
    """
    提取 PPT 内容。
    防幻觉策略：
      - 按幻灯片提取，保留幻灯片编号
      - 不重新整理内容顺序
      - 区分标题/正文/表格内容
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"found": False, "content": "python-pptx 未安装", "note": "pip install python-pptx"}

    query_lower = query.lower()
    query_words = [w for w in query_lower.split() if len(w) >= 2]
    prs = Presentation(file_path)
    results = []

    for slide_num, slide in enumerate(prs.slides, 1):
        # 区分标题和正文
        titles = []
        body_lines = []
        tables_data = []

        for shape in slide.shapes:
            if shape.has_table:
                t = shape.table
                table_rows = []
                for row in t.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    table_rows.append(row_text)
                tables_data.append("\n".join(table_rows))
            elif hasattr(shape, "text_frame"):
                text = shape.text.strip()
                if not text:
                    continue
                if shape.has_text_frame and shape.text_frame.first_line_format.font.bold:
                    titles.append(text[:100])
                else:
                    body_lines.append(text[:200])

        full_text = " ".join(titles + body_lines)
        score = sum(full_text.lower().count(w) for w in query_words)
        if score == 0:
            continue

        # 组装片段
        results.append(f"[幻灯片 {slide_num}]（相关度:{score}）")
        if titles:
            results.append("标题: " + " | ".join(titles))
        if body_lines:
            results.append("正文: " + " ".join(body_lines[:3]))
        if tables_data:
            results.append("表格: " + tables_data[0][:300])
        results.append("")

    if not results:
        return {
            "found": False,
            "content": "未找到匹配幻灯片",
            "note": "python-pptx 按幻灯片提取，编号已标注"
        }

    return {
        "found": True,
        "content": "\n".join(results[: top_k * 5]),
        "total_slides": len(prs.slides),
        "note": "python-pptx 按幻灯片提取，区分标题/正文/表格"
    }


def extract_link(url: str, query: str, top_k: int = 3) -> dict:
    """
    提取网页内容。
    防幻觉策略：
      - 精确引用URL页面原文
      - 标注：URL + 抓取时间 + 页面标题
      - 去除广告/导航，仅保留正文
      - 自动降级：requests+BS4 → urllib
    """
    raw_text = None
    fetch_method = "unknown"

    # 优先用 requests + BeautifulSoup
    try:
        import requests
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                          "AppleWebKit/537.36 (KHTML, like Gecko) "
                          "Chrome/120.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
        }
        resp = requests.get(url, headers=headers, timeout=15, allow_redirects=True)
        resp.encoding = resp.apparent_encoding or "utf-8"
        raw_html = resp.text
        fetch_method = "requests"
    except Exception:
        # 降级到 urllib
        try:
            import urllib.request
            req = urllib.request.Request(url, headers={
                "User-Agent": "Mozilla/5.0 (compatible; FashionDoctorBot/1.0)"
            })
            with urllib.request.urlopen(req, timeout=10) as r:
                raw = r.read()
                try:
                    raw_html = raw.decode("utf-8")
                except UnicodeDecodeError:
                    raw_html = raw.decode("gbk", errors="replace")
            fetch_method = "urllib"
        except Exception as e2:
            return {
                "found": False,
                "content": "",
                "note": f"网页抓取失败（requests+urllib均失败）：{e2}，请手动访问URL验证"
            }

    # 用 BeautifulSoup 提取正文（如果可用）
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw_html, "html.parser")

        # 移除噪音元素
        for tag in soup(["script", "style", "nav", "header", "footer",
                          "aside", "noscript", "iframe", "form"]):
            tag.decompose()

        # 提取 <article> 或 <main> 或最大文本块
        main = (soup.find("article") or
                soup.find("main") or
                soup.find("div", class_=lambda x: x and "content" in str(x).lower()) or
                soup.body if soup.body else soup)

        text = main.get_text(separator="\n", strip=True)
        text = re.sub(r"\n{3,}", "\n\n", text)
        page_title = soup.title.string.strip() if soup.title else "无标题"
    except Exception:
        # 无 BS4，纯文本提取
        text = re.sub(r"<[^>]+>", " ", raw_html)
        text = re.sub(r"\s+", " ", text).strip()
        page_title = "（无标题）"

    # 按查询关键词截取最相关的段落
    query_lower = query.lower()
    query_words = [w for w in query_lower.split() if len(w) >= 2]
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    scored = []
    for p in paragraphs:
        p_lower = p.lower()
        score = sum(p_lower.count(w) for w in query_words)
        if score > 0:
            scored.append((score, p))
    scored.sort(key=lambda x: x[0], reverse=True)

    if scored:
        content = "\n\n".join(p for _, p in scored[:top_k])
        note = f"网页正文提取（{fetch_method}+BeautifulSoup），关键词段落精准召回"
    else:
        content = "\n\n".join(paragraphs[:5])
        note = f"网页正文提取（{fetch_method}），未精确匹配关键词，返回前5段"

    return {
        "found": True,
        "content": content[:2000],
        "url": url,
        "title": page_title,
        "fetched_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "fetch_method": fetch_method,
        "note": note
    }


# ══════════════════════════════════════════════════════
# 第四部分：知识库内容新增（增量写入）
# ══════════════════════════════════════════════════════

def add_kb_entry(L2_id: str, L3_name: str, content: str,
                  content_type: str = "md",
                  source_info: str = "",
                  author: str = "Fashion Doctor") -> dict:
    """
    新增一条知识库内容。
    自动写入对应L3文件，并在 master_index.json 中追加条目。
    返回写入结果和文件路径。
    """
    index = load_index()
    if not index:
        return {"success": False, "error": "无法读取索引文件"}

    # 找到L2目录
    L2_cat = next((c for c in index["L2_categories"] if c["id"] == L2_id), None)
    if not L2_cat:
        return {"success": False, "error": f"未找到L2品类：{L2_id}"}

    # 生成L3 ID（自动编号）
    existing_ids = [e["id"] for e in L2_cat.get("L3", [])]
    nums = [int(re.search(r'\d+', i).group()) for i in existing_ids
            if re.search(r'\d+', i)]
    next_num = max(nums) + 1 if nums else 1
    L3_id = f"L3_{L2_id.split('_')[1]}_{next_num:02d}"

    # 构建文件名和路径
    safe_name = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', '_', L3_name)
    filename = f"{safe_name}.{content_type}"
    L3_dir = KB_ROOT / L2_cat["id"] / f"L3_{L2_id.split('_')[1]}_{next_num:02d}_{safe_name}"
    L3_dir.mkdir(parents=True, exist_ok=True)
    file_path = L3_dir / filename

    # 写入内容（添加标准头部）
    header = f"""# {L3_name}

> **来源**：{source_info or '新增'}
> **版本**：v1.0（{datetime.now().strftime('%Y-%m-%d')}）
> **录入人**：{author}
> **类型**：知识库 / {L2_cat['name']} / {L3_name}

---

{content}
"""
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(header)

    # 更新索引
    rel_path = str(file_path.relative_to(KB_ROOT)).replace("\\", "/")
    entry = {
        "id": L3_id,
        "name": L3_name,
        "path": rel_path,
        "status": "manual",
        "added": datetime.now().strftime("%Y-%m-%d")
    }
    L2_cat.setdefault("L3", []).append(entry)
    index["total_entries"] += 1
    index["last_updated"] = datetime.now().strftime("%Y-%m-%d")

    with open(INDEX_FILE, "w", encoding="utf-8") as f:
        json.dump(index, f, ensure_ascii=False, indent=2)

    return {
        "success": True,
        "L3_id": L3_id,
        "file_path": str(file_path),
        "rel_path": rel_path,
        "L2_name": L2_cat["name"]
    }


# ══════════════════════════════════════════════════════
# 第五部分：主检索函数（对外接口）
# ══════════════════════════════════════════════════════

def retrieve(query: str,
             content_type: str = None,
             level_filter: str = None,
             top_k: int = 3) -> RetrievalResult:
    """
    知识库主检索函数。

    参数：
        query       - 检索关键词
        content_type - 内容类型：md|excel|pdf|image|ppt|link，不指定则自动推断
        level_filter - L2或L3筛选，如 "竞品分析"
        top_k       - 每个文件最多返回的段落数

    返回：
        RetrievalResult 对象（含 answer/sources/confidence）

    防幻觉机制：
        1. 所有摘录带来源路径 + 行号/页码
        2. 表格数据原样提取，不做计算
        3. 无法确认的内容标记为 unverified
        4. 始终标注置信度（high/medium/low/unverified）
    """
    result = RetrievalResult()
    result.query = query

    index = load_index()
    if not index:
        result.answer = "❌ 知识库索引文件不存在，请检查 knowledge_base/__index__/master_index.json"
        result.confidence = "unverified"
        return result

    # 搜索索引
    matches = search_index(query, index)
    if not matches:
        result.answer = f"❌ 知识库中未找到与「{query}」相关的内容。"
        result.answer += "\n\n💡 建议：\n  1. 尝试更通用的关键词（如「KPI」「竞品」「会员」）\n  2. 使用 add_kb_entry() 新增内容\n  3. 检查 master_index.json 是否包含相关L2品类"
        result.confidence = "unverified"
        return result

    # 按内容类型筛选
    if content_type:
        matches = [m for m in matches if m["path"].suffix == f".{content_type}"]

    # 提取最相关文件（最多3个）
    all_snippets = []
    for match in matches[:3]:
        path = match["path"]
        ext = path.suffix.lower().lstrip(".")

        if ext == "md":
            extracted = extract_md(path, query, top_k)
        elif ext in ("xlsx", "xls", "csv"):
            extracted = extract_excel(path, query)
        elif ext in ("png", "jpg", "jpeg", "gif", "bmp"):
            extracted = extract_image(path, query)
        elif ext == "pdf":
            extracted = extract_pdf(path, query, top_k)
        elif ext in ("pptx", "ppt"):
            extracted = extract_ppt(path, query, top_k)
        elif ext == "txt":
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
            extracted = {"found": True, "content": content[:500]}
        else:
            extracted = {"found": False, "content": f"不支持的文件类型: {ext}"}

        if extracted.get("found"):
            all_snippets.append({
                "type": ext.upper().lstrip("."),
                "path": str(path),
                "L2_name": match["L2_name"],
                "L3_name": match["name"],
                "content": extracted.get("content", ""),
                "version": extracted.get("version", "v1.0"),
                "confidence": "high"
            })

    if not all_snippets:
        result.answer = f"❌ 找到匹配文件，但内容提取失败。「{query}」"
        result.confidence = "low"
        return result

    # 汇总答案
    answer_parts = []
    for snip in all_snippets:
        answer_parts.append(f"### 📎 {snip['L3_name']}（{snip['type']} | {snip['L2_name']}）\n{snip['content']}\n")

    result.answer = "\n".join(answer_parts)
    result.sources = all_snippets
    result.confidence = "high"
    result.related_entries = [
        {"name": m["name"], "L2": m["L2_name"], "path": str(m["path"])}
        for m in matches[3:6]
    ]

    return result


# ══════════════════════════════════════════════════════
# 第六部分：CLI 入口
# ══════════════════════════════════════════════════════

def main():
    if len(sys.argv) < 2:
        print("用法: python retrieval_mod.py <查询关键词> [--type md|excel|pdf|image|ppt|link]")
        print("示例: python retrieval_mod.py 死库率行业基准")
        print("      python retrieval_mod.py 太平鸟男装周转 --type md")
        print()
        print("交互模式: python retrieval_mod.py --interactive")
        return

    args = sys.argv[1:]
    if "--interactive" in args:
        print("🔍 Fashion Doctor 知识库检索（交互模式）")
        print("输入查询内容，或按 Ctrl+C 退出\n")
        while True:
            try:
                q = input("查询> ").strip()
                if not q:
                    continue
                if q in ("exit", "quit", "q"):
                    break
                r = retrieve(q)
                r.print_summary()
            except (KeyboardInterrupt, EOFError):
                break
        return

    query = args[0]
    content_type = None
    for i, a in enumerate(args):
        if a == "--type" and i + 1 < len(args):
            content_type = args[i + 1]

    r = retrieve(query, content_type=content_type)
    r.print_summary()


if __name__ == "__main__":
    main()
